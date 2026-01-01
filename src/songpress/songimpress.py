###############################################################
# Name:             songimpress.py
# Purpose:     Generate PowerPoint presentation from songs
# Author:         Luca Allulli (webmaster@roma21.it)
# Created:     2019-02-02
# Copyright: Luca Allulli (https://www.skeed.it/songpress)
# License:     GNU GPL v2
##############################################################

import sys

from pptx import Presentation


class SongPresentation:
    def __init__(self, template, out_filename):
        self.pres = Presentation(template)
        self.layout = self.pres.slide_layouts[3]
        self.prev = None
        self.out_filename = out_filename

    def _add_slide(self, cur, next):
        slide = self.pres.slides.add_slide(self.layout)
        ph = list(slide.placeholders)
        t1 = ph[1]
        t2 = ph[2]
        t1.text = cur
        t2.text = next

    def _add_empty_slide(self):
        self.pres.slides.add_slide(self.layout)

    def add_line(self, line):
        if self.prev is not None:
            self._add_slide(self.prev, line)
        self.prev = line

    def end_song(self):
        if self.prev is not None:
            self._add_slide(self.prev, '')
        self._add_empty_slide()
        self.prev = None

    def close(self):
        if self.prev is not None:
            self._add_slide(self.prev, '')
        self.pres.save(self.out_filename)

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        self.close()


def to_presentation(lines, output_file, template_file):
    with SongPresentation(template_file, output_file) as c:
        for line in lines:
            line = line.replace('\n', '').strip()
            if line == '---':
                c.end_song()
            elif line != '':
                c.add_line(line)


if __name__ == '__main__':
    if len(sys.argv) != 4:
        print("Usage: songimpress.py SONG_FILE OUTPUT_PRESENTATION TEMPLATE_PRESENTATION")
        exit(1)

    with open(sys.argv[1]) as f:
        to_presentation(f, sys.argv[2], sys.argv[3])
